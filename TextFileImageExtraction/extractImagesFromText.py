from docx import Document
import PyPDF2
from pptx import Presentation
import os
from PIL import Image
import zlib
import io
import os
import csv

# Extract all Images from the DOCX File

def extract_images_from_docx(filename):
    doc = Document(filename)
    image_counter = 0
    # Create a new folder with the same name as the file
    output_dir = os.path.join("New DATA/TextToImages/DOCTextImages", os.path.basename(filename))
    os.makedirs(output_dir, exist_ok=True)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image = rel._target
            image_data = rel.target_part.blob
            
            image_filename = f"image_{image_counter}.png"
            image_path = os.path.join(output_dir, image_filename)

            with open(image_path, "wb") as f:
                f.write(image_data)

            image_counter += 1

    print(f"Extracted {image_counter} image(s) from the document.")
    return output_dir

# Extract all Images from the PDF Files


def sanitize_filename(filename):
    # Replace characters not allowed in file names
    invalid_chars = r'\/:*?"<>|'
    for char in invalid_chars:
        filename = filename.replace(char, "_")
    return filename

def extract_images_from_pdf(pdf_file):
    filename=os.path.basename(pdf_file)
    output_dir = os.path.join("New DATA/TextToImages/PDFTextImages", filename)
    os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist

    with open(pdf_file, "rb") as file:
        pdf = PyPDF2.PdfReader(file)
        
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            if "/XObject" in page["/Resources"]:
                x_object = page["/Resources"]["/XObject"].get_object()
                for obj in x_object:
                    if x_object[obj]["/Subtype"] == "/Image":
                        image = x_object[obj]

                        if "/Filter" in image:
                            filters = image["/Filter"]

                            if "/FlateDecode" in filters:
                                # Decompress the FlateDecode image data
                                image_data = zlib.decompress(image._data)
                            elif "/DCTDecode" in filters:
                                # Image is already in JPEG format, no decoding needed
                                image_data = image._data
                            else:
                                # Unsupported image encoding
                                print(f"Skipping image {obj}: Unsupported image encoding")
                                continue
                        else:
                            # No encoding specified, treat as FlateDecode
                            image_data = zlib.decompress(image._data)

                        # Create an image object from the image data
                        image_stream = io.BytesIO(image_data)
                        try:
                            image_obj = Image.open(image_stream)
                        except IOError:
                            # print(f"Skipping image {obj}: Error in identifying image format")
                            continue

                        # Save the image as PNG
                        image_name = f"image_{page_num}_{sanitize_filename(obj)}.png"
                        image_path = os.path.join(output_dir, image_name)
                        image_obj.save(image_path, "PNG")
    return output_dir

def extract_images_from_ppt(filename):
    prs = Presentation(filename)
    output_dir = os.path.join("New DATA/TextToImages/PPTXextImages", os.path.basename(filename))
    os.makedirs(output_dir, exist_ok=True)
    image_counter = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is an image
                image = shape.image
                image_bytes = image.blob
                with open(f"image_{image_counter}.png", "wb") as f:
                    f.write(image_bytes)

                image_counter += 1

    print(f"Extracted {image_counter} image(s) from the PowerPoint.")
    return output_dir



def loadFileImagesToCSV(folder_path):
    print("TTTTTTTTTTTTTT:",folder_path)
    # create a CSV file and open it in write mode
    with open('TestAndResults\FilesImagesData.csv', 'a', newline='') as csvfile:
        # create a CSV writer
        csvwriter = csv.writer(csvfile)
        # write the header row
        csvwriter.writerow(['Folder Name', 'File Name', 'File Path'])
        
        # use os.walk() to get all the files and folders in the specified path
        for root, dirs, files in os.walk(folder_path):
            # loop through the list of files in the current directory
            for file in files:
                # construct the full path to the file
                file_path = os.path.join(root, file)
                
                # extract the last directory name from the file path
                folder_name = os.path.basename(os.path.dirname(file_path))
                
                # write the data to the CSV file
                csvwriter.writerow([folder_name, file, file_path])
                
                print("Root Path =", root)
                print("File Name =", file)
                print("File Path =", file_path)
                print("Folder Name =", folder_name)






#loadFileImages("New DATA\TextToImages")           
# # Replace "your_ppt.pptx" with the path to your PowerPoint file
# ppt_file = "your_ppt.pptx"
# extract_images_from_ppt(ppt_file)
# # Usage example
# extract_images_from_pdf("New DATA\documents\AI_Manual.pdf")
# extract_images_from_docx("New DATA\\documents\\Material.docx")
# extract_images_from_ppt()